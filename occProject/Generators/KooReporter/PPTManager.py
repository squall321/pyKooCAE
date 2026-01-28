import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Cm
import matplotlib.pyplot as plt
import re

class SlideContent:
    def __init__(self):
        self.slide = None
    def add_to_presentation(self, presentation):
        raise NotImplementedError("서브클래스에서 구현하세요.")

    def make_text_box(self, x, y, width, height, text, font_size = 18):
        textbox = self.slide.shapes.add_textbox(Inches(x), Inches(y), Inches(width), Inches(height))
        tf = textbox.text_frame
        for i, line in enumerate(text):
            p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
            '''if i == 0:
                p.text = line  # 일반 글머리
                p.level = 0
            else:
                p.text = f"- {line}"  # 대시 붙이기
                p.level = 1'''
            if i == 0:
                p.text = f"■ {line}"  # 첫 줄: 원형 글머리
                p.font.bold = True  # ← 볼드체 설정

            else:
                p.text = f" - {line}"  # 이후 줄: 대시
            p.level = 0
            p.font.size = Pt(font_size)
            
            
    def make_image_box(self, x, y, width, height, image_path):
        image_x = Inches(x)
        image_y = Inches(y)
        image_width = Inches(width)
        image_height = Inches(height)
        if os.path.exists(image_path):
            self.slide.shapes.add_picture(image_path, image_x, image_y, height=image_height)
        else:
            # 사각형 대체
            shape = self.slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, image_x, image_y, Inches(4), image_height)
            shape.text = "이미지 없음"
            
    # 1. 수식 이미지를 생성하는 함수
    def render_formula_to_image(self, formula: str, filename: str = "equation.png"):
        curPath = os.path.dirname(os.path.abspath(__file__))
        plt.figure(figsize=(2, 1))
        plt.text(0.5, 0.5, f"${formula}$", fontsize=30, ha='center', va='center')
        plt.axis('off')
        curPath = os.path.join(curPath, filename)
        plt.savefig(curPath, bbox_inches='tight', dpi=300, transparent=True)
        plt.close()
        return curPath
    
    def make_text_with_equation(self, x, y, lines, font_size=18):
        for i, line in enumerate(lines):
            offset_y = y + i * 0.4  # 줄 간격
            if i == 0:
                bold = True
                header = "■ "  # 첫 줄은 원형 글머리
            else:
                bold = False
                header = " - "  # 이후 줄은 대시
             
            self.insert_text_and_formula_line(line, x, offset_y, font_size, bold, header)
            
    def insert_text_and_formula_line(self, text, x, y, font_size=18, bold=False, header=" - "):
        parts = re.split(r'(\$.*?\$)', text)
        cur_x = x
        ii = 0

        for part in parts:
            if part.startswith('$') and part.endswith('$'):
                formula = part[1:-1]
                img_path = self.render_formula_to_image(formula)
                pic = self.slide.shapes.add_picture(img_path, Inches(cur_x), Inches(y + 0.04), height=Inches(0.4))
                width_inches = pic.width.inches
                cur_x += width_inches + 0.0
            else:
                # 🔧 텍스트 길이에 따라 박스 폭 추정
                if ii == 0:
                    estimated_width = (len(header)+len(part)) * 0.20  # 첫 줄은 더 넓게
                else:
                    estimated_width = len(part) * 0.20
                
                
                textbox = self.slide.shapes.add_textbox(Inches(cur_x), Inches(y), Inches(estimated_width), Inches(0.4))
                tf = textbox.text_frame
                p = tf.paragraphs[0]
                if ii == 0:
                    p.text = f"{header}{part}"
                else:
                    p.text = part
                p.font.size = Pt(font_size)
                p.font.bold = bold
                cur_x += estimated_width + 0.02
            ii += 1
                
class SingleColumnSlide(SlideContent):
    def __init__(self, title, description_lines, image_path):
        super().__init__()
        self.title = title
        self.description_lines = description_lines[:3]
        self.image_path = image_path

    def add_to_presentation(self, presentation):
        self.slide = presentation.slides.add_slide(presentation.slide_layouts[6])  # 빈 슬라이드

        # 제목
        title_box = self.slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.6))

        title_frame = title_box.text_frame
        title_frame.text = self.title
        title_frame.paragraphs[0].font.size = Pt(30)
        title_frame.paragraphs[0].font.bold = True
        title_frame.paragraphs[0].alignment = PP_ALIGN.LEFT

        
        # 본문 설명
        #self.make_text_box(0.5, 1.0, 12, 1.0, self.description_lines, font_size=18)
        self.make_text_with_equation(0.5, 1.0, self.description_lines)

        # 이미지 (하단 중앙)
        self.make_image_box(4.2, 2.5, 5, 4, self.image_path)
       


class TwoColumnSlide(SlideContent):
    def __init__(self, title, left_desc, left_img, right_desc, right_img):
        super().__init__()
        self.title = title
        self.left_desc = left_desc[:3]
        self.left_img = left_img
        self.right_desc = right_desc[:3]
        self.right_img = right_img

    def add_to_presentation(self, presentation):
        self.slide = presentation.slides.add_slide(presentation.slide_layouts[6])  # 빈 슬라이드

        # 제목
        title_box = self.slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.6))

        title_frame = title_box.text_frame
        title_frame.text = self.title
        title_frame.paragraphs[0].font.size = Pt(30)
        title_frame.paragraphs[0].font.bold = True
        title_frame.paragraphs[0].alignment = PP_ALIGN.LEFT

        # 왼쪽 설명
        self.make_text_box(0.5, 1.0, 6, 1.0, self.left_desc, font_size=18)
        # 오른쪽 설명
        self.make_text_box(7.0, 1.0, 6, 1.0, self.right_desc, font_size=18)
        
        # 왼쪽 이미지
        self.make_image_box(0.5, 2.2, 5.5, 4, self.left_img)
        
        # 오른쪽 이미지
        self.make_image_box(7.0, 2.2, 5.5, 4, self.right_img)


class PPTManager:
    def __init__(self):
        self.contents = []
        self.filename = "output.pptx"
        self.presentation = Presentation()
        self.presentation.slide_width = Inches(13.33)
        self.presentation.slide_height = Inches(7.5)

    def SetFileName(self, filename):
        self.filename = filename

    def AddContent(self, content: SlideContent):
        self.contents.append(content)

    def AddtoPresentation(self):
        for content in self.contents:
            content.add_to_presentation(self.presentation)

    def Print(self):
        self.presentation.save(self.filename)



if __name__ == "__main__":
    manager = PPTManager()
    manager.SetFileName("my_slides.pptx")

    # 1열 슬라이드
    firstSlide = SingleColumnSlide(
        "제품 설명",
        ["이 제품은 매우 강력합니다.", "다양한 환경에서 사용할 수 있습니다.", "경제적이기도 합니다."],
        "left_image.png"
    )
    manager.AddContent(firstSlide)
    
    secondSlide = TwoColumnSlide(
        "비교 분석",
        ["A 제품은 빠릅니다.", "디자인이 세련됨", "저렴함"], "left_image.png",
        ["B 제품은 안정적입니다.", "오랜 수명", "높은 성능"], "right_image.png"
    )

    # 2열 슬라이드
    manager.AddContent(secondSlide)
    
    thirdSlide = SingleColumnSlide(
        "수식 테스트 페이지",
        [
            "에너지는 질량과 빛의 속도로 표현되며 $E=mc^{2}$ 이다.",
            "삼각형의 피타고라스 정리는 $a^2 + b^2 = c^2$ 로 쓴다.",
            "면적은 $A = \\pi r^2$ 으로 계산된다."
        ],
        "none.png"
    )
    
    manager.AddContent(thirdSlide)
    manager.AddtoPresentation()
    manager.Print()

